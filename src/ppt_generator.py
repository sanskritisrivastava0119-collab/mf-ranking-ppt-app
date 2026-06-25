from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor


APP_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = APP_ROOT / "assets" / "ranking-template.pptx"
MAX_SCHEMES = 200
ROWS_PER_SLIDE = 8
BLACK = RGBColor(0, 0, 0)


def _cell_text(cell) -> str:
    return "\n".join(paragraph.text for paragraph in cell.text_frame.paragraphs).strip()


def _set_cell_text_preserving_style(cell, value: str) -> None:
    text_frame = cell.text_frame
    if not text_frame.paragraphs:
        text_frame.text = value
        _set_cell_text_black(cell)
        return

    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = value
        for extra_run in paragraph.runs[1:]:
            extra_run.text = ""
    else:
        paragraph.text = value

    for extra_paragraph in text_frame.paragraphs[1:]:
        extra_paragraph.text = ""
    _set_cell_text_black(cell)


def _set_cell_text_black(cell) -> None:
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = BLACK


def _ranking_table(slide):
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        first_row = [_cell_text(cell) for cell in table.rows[0].cells]
        if "Scheme" in first_row and "Category" in first_row:
            return table
    raise ValueError("Ranking table not found in the PowerPoint template.")


def _ensure_body_rows(table, body_rows: int) -> None:
    while len(table.rows) - 2 < body_rows:
        new_row = deepcopy(table._tbl.tr_lst[-1])
        table._tbl.append(new_row)


def _remove_slide(presentation: Presentation, index: int) -> None:
    slide_id = presentation.slides._sldIdLst[index]
    relationship_id = slide_id.rId
    presentation.part.drop_rel(relationship_id)
    presentation.slides._sldIdLst.remove(slide_id)


def _duplicate_slide(presentation: Presentation, source_index: int) -> None:
    source = presentation.slides[source_index]
    new_slide = presentation.slides.add_slide(presentation.slide_layouts[0])

    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)

    for shape in source.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element), "p:extLst"
        )

    for relationship in source.part.rels.values():
        if "notesSlide" in relationship.reltype or "slideLayout" in relationship.reltype:
            continue
        new_slide.part.rels.add_relationship(
            relationship.reltype,
            relationship._target,
            relationship.rId,
        )


def build_ranking_deck(rows: Iterable[dict[str, str]]) -> bytes:
    records = list(rows)
    if not records:
        raise ValueError("At least one completed ranking row is required.")
    if len(records) > MAX_SCHEMES:
        raise ValueError(f"The template supports a maximum of {MAX_SCHEMES} schemes.")
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Missing template: {TEMPLATE_PATH}")

    presentation = Presentation(str(TEMPLATE_PATH))
    required_slides = (len(records) + ROWS_PER_SLIDE - 1) // ROWS_PER_SLIDE

    while len(presentation.slides) < required_slides:
        _duplicate_slide(presentation, min(len(presentation.slides) - 1, 2))

    for slide_index, slide in enumerate(presentation.slides):
        table = _ranking_table(slide)
        _ensure_body_rows(table, ROWS_PER_SLIDE)
        page_rows = records[
            slide_index * ROWS_PER_SLIDE : (slide_index + 1) * ROWS_PER_SLIDE
        ]

        for row_offset in range(len(table.rows) - 2):
            table_row = table.rows[row_offset + 2]
            if row_offset < len(page_rows):
                record = page_rows[row_offset]
                values = [
                    record["Scheme"],
                    record["Category"],
                    record["1Y"],
                    record["3Y"],
                    record["5Y"],
                ]
            else:
                values = ["", "", "", "", ""]

            for cell, value in zip(table_row.cells, values):
                _set_cell_text_preserving_style(cell, str(value))

    while len(presentation.slides) > required_slides:
        _remove_slide(presentation, len(presentation.slides) - 1)

    output = __import__("io").BytesIO()
    presentation.save(output)
    return output.getvalue()
