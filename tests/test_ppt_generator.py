from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor

from src.ppt_generator import build_ranking_deck


def test_build_deck_removes_unused_slides():
    data = [
        {
            "Scheme": "Example Fund",
            "Category": "Flexi Cap Fund",
            "1Y": "1/10",
            "3Y": "2/9",
            "5Y": "3/8",
        }
    ]

    result = build_ranking_deck(data)
    presentation = Presentation(BytesIO(result))

    assert len(presentation.slides) == 1
    tables = [shape.table for shape in presentation.slides[0].shapes if shape.has_table]
    assert "Example Fund" in tables[0].cell(2, 0).text


def test_build_deck_supports_two_hundred_rows_across_generated_slides():
    data = [
        {
            "Scheme": f"Example Fund {index}",
            "Category": "Flexi Cap",
            "1Y": "1/10",
            "3Y": "2/9",
            "5Y": "3/8",
        }
        for index in range(1, 201)
    ]

    result = build_ranking_deck(data)
    presentation = Presentation(BytesIO(result))

    assert len(presentation.slides) == 25
    last_table = [
        shape.table for shape in presentation.slides[24].shapes if shape.has_table
    ][0]
    assert "Example Fund 200" in last_table.cell(9, 0).text


def test_build_deck_writes_all_body_values_in_black():
    data = [
        {
            "Scheme": "Example Fund",
            "Category": "Flexi Cap",
            "1Y": "1/10",
            "3Y": "2/9",
            "5Y": "3/8",
        }
    ]

    result = build_ranking_deck(data)
    presentation = Presentation(BytesIO(result))
    table = [shape.table for shape in presentation.slides[0].shapes if shape.has_table][0]

    for cell in table.rows[2].cells:
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    assert run.font.color.rgb == RGBColor(0, 0, 0)
